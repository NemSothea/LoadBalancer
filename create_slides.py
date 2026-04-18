#!/usr/bin/env python3
"""
Generate a polished PowerPoint presentation:
  Load Balancer Setup with Nginx on AlmaLinux
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
BG_CARD      = RGBColor(0x16, 0x21, 0x3E)   # card background
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)   # cyan accent
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)   # green accent
ACCENT_PURPLE= RGBColor(0xA7, 0x8B, 0xFA)   # purple accent
ACCENT_ORANGE= RGBColor(0xFB, 0x92, 0x3C)   # orange accent
ACCENT_PINK  = RGBColor(0xF4, 0x72, 0xB6)   # pink accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
CODE_BG      = RGBColor(0x1E, 0x29, 0x3B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

IMG_DIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def add_bg(slide, color=BG_DARK):
    """Fill the slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=16,
                      color=WHITE, font_name="Calibri", alignment=PP_ALIGN.LEFT,
                      bold=False, line_spacing=1.5):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rich_text_box(slide, left, top, width, height, runs_list,
                      alignment=PP_ALIGN.LEFT):
    """
    runs_list = [
        [ (text, size, color, bold, font_name), ... ],  # paragraph 1
        [ (text, size, color, bold, font_name), ... ],  # paragraph 2
    ]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for pi, para_runs in enumerate(runs_list):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(6)
        for ri, (text, size, color, bold, font_name) in enumerate(para_runs):
            if ri == 0:
                run = p.runs[0] if p.runs else p.add_run()
                run.text = text
            else:
                run = p.add_run()
                run.text = text
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = font_name
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Draw a thin accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bullet_card(slide, left, top, width, height, title, bullets,
                    title_color=ACCENT_BLUE, bullet_color=LIGHT_GRAY,
                    card_color=BG_CARD):
    """Add a card with a title and bullet list."""
    add_shape(slide, left, top, width, height, card_color, 0.05)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2),
                 width - Inches(0.6), Inches(0.5),
                 title, font_size=20, color=title_color, bold=True)
    add_accent_line(slide, left + Inches(0.3), top + Inches(0.65),
                    Inches(1.5), title_color)
    lines = []
    for b in bullets:
        lines.append(b)
    add_multiline_box(slide, left + Inches(0.3), top + Inches(0.8),
                      width - Inches(0.6), height - Inches(1.0),
                      lines, font_size=14, color=bullet_color,
                      line_spacing=1.6)


def add_code_block(slide, left, top, width, height, code_text):
    """Add a styled code block."""
    add_shape(slide, left, top, width, height, CODE_BG, 0.03)
    add_multiline_box(slide, left + Inches(0.25), top + Inches(0.15),
                      width - Inches(0.5), height - Inches(0.3),
                      code_text.split("\n"), font_size=12,
                      color=ACCENT_GREEN, font_name="Courier New",
                      line_spacing=1.4)


def add_slide_number(slide, num, total):
    """Add a slide number in bottom-right."""
    add_text_box(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.5),
                 Inches(1.2), Inches(0.4),
                 f"{num} / {total}", font_size=11, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_section_badge(slide, text, color=ACCENT_BLUE):
    """Small badge at top-left."""
    shape = add_shape(slide, Inches(0.6), Inches(0.4), Inches(2.2), Inches(0.4), color, 0.15)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.color.rgb = BG_DARK
    run.font.bold = True
    run.font.name = "Calibri"


TOTAL_SLIDES = 16

# ===================================================================
# SLIDE 1 — TITLE
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# Decorative shapes
add_shape(slide, Inches(-1), Inches(-1), Inches(6), Inches(6),
          RGBColor(0x38, 0xBD, 0xF8), 0.5).fill.fore_color.rgb = RGBColor(0x14, 0x28, 0x50)
add_shape(slide, Inches(9), Inches(3), Inches(6), Inches(6),
          RGBColor(0x38, 0xBD, 0xF8), 0.5).fill.fore_color.rgb = RGBColor(0x14, 0x28, 0x50)

add_accent_line(slide, Inches(1.5), Inches(2.4), Inches(3), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2),
             "Load Balancer Setup", font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
             "with Nginx on AlmaLinux", font_size=36, color=ACCENT_BLUE, bold=False)

add_accent_line(slide, Inches(1.5), Inches(4.7), Inches(5), ACCENT_PURPLE)

add_text_box(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.5),
             "High-Availability Web Infrastructure  |  Step-by-Step Guide",
             font_size=18, color=LIGHT_GRAY)
add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10), Inches(0.5),
             "April 2026  |  AlmaLinux 10  |  Nginx 1.26  |  PHP-FPM",
             font_size=14, color=MID_GRAY)

add_slide_number(slide, 1, TOTAL_SLIDES)

# ===================================================================
# SLIDE 2 — TABLE OF CONTENTS
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "OVERVIEW")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.7),
             "Table of Contents", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_BLUE)

toc_items = [
    ("01", "Project Overview", ACCENT_BLUE),
    ("02", "Why Use a Load Balancer?", ACCENT_GREEN),
    ("03", "Architecture & Process Flow", ACCENT_PURPLE),
    ("04", "Environment Details", ACCENT_ORANGE),
    ("05", "Phase 1 : VM1 - Websites + Dashboard", ACCENT_BLUE),
    ("06", "Phase 1 : Configuration & SELinux", ACCENT_GREEN),
    ("07", "Phase 2 : VM2 - Load Balancer Setup", ACCENT_PURPLE),
    ("08", "Demo : Site A & Site B", ACCENT_ORANGE),
    ("09", "Demo : Health Dashboard", ACCENT_PINK),
    ("10", "Demo : Nginx Load Balancer in Action", ACCENT_BLUE),
    ("11", "Verification & Testing", ACCENT_GREEN),
    ("12", "Challenges & Solutions", ACCENT_PURPLE),
    ("13", "Extending the Infrastructure", ACCENT_ORANGE),
    ("14", "Conclusion", ACCENT_PINK),
]

col1 = toc_items[:7]
col2 = toc_items[7:]

for i, (num, title, color) in enumerate(col1):
    y = Inches(2.3) + Inches(i * 0.6)
    add_shape(slide, Inches(0.8), y, Inches(0.55), Inches(0.4), color, 0.15)
    add_text_box(slide, Inches(0.8), y - Pt(1), Inches(0.55), Inches(0.4),
                 num, font_size=13, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), y, Inches(5), Inches(0.4),
                 title, font_size=15, color=LIGHT_GRAY)

for i, (num, title, color) in enumerate(col2):
    y = Inches(2.3) + Inches(i * 0.6)
    add_shape(slide, Inches(7.0), y, Inches(0.55), Inches(0.4), color, 0.15)
    add_text_box(slide, Inches(7.0), y - Pt(1), Inches(0.55), Inches(0.4),
                 num, font_size=13, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.7), y, Inches(5), Inches(0.4),
                 title, font_size=15, color=LIGHT_GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)

# ===================================================================
# SLIDE 3 — PROJECT OVERVIEW
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "01  PROJECT OVERVIEW")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Project Overview", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.6),
             "Two virtual machines working together as a high-availability web infrastructure",
             font_size=18, color=LIGHT_GRAY)

# VM1 Card
add_bullet_card(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(3.5),
                "VM1  -  Web Server  (192.168.1.8)",
                [
                    "  Site A on port 8081 (Primary Web App)",
                    "  Site B on port 8082 (Secondary Web App)",
                    "  Health Dashboard on port 8080",
                    "  PHP-FPM for dynamic content",
                    "  Real-time uptime monitoring",
                ],
                title_color=ACCENT_BLUE)

# VM2 Card
add_bullet_card(slide, Inches(7.0), Inches(3.0), Inches(5.5), Inches(3.5),
                "VM2  -  Load Balancer  (192.168.1.20)",
                [
                    "  Nginx reverse proxy on port 80",
                    "  Round-robin traffic distribution",
                    "  Proxies to VM1:8081 and VM1:8082",
                    "  Automatic failover on backend failure",
                    "  Single public entry point for users",
                ],
                title_color=ACCENT_PURPLE)

add_slide_number(slide, 3, TOTAL_SLIDES)

# ===================================================================
# SLIDE 4 — WHY LOAD BALANCER?
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "02  WHY LOAD BALANCER?")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Why Use a Load Balancer?", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_GREEN)

reasons = [
    ("High Availability", "If one backend fails, traffic is automatically\nredirected to the healthy one.", ACCENT_BLUE),
    ("Scalability", "Add more backend servers without changing\nthe public endpoint.", ACCENT_GREEN),
    ("Traffic Distribution", "Prevents overload on a single server using\nround-robin, least connections, etc.", ACCENT_PURPLE),
    ("Zero-Downtime Maintenance", "Take one backend offline while the\nother continues serving requests.", ACCENT_ORANGE),
    ("Simplified SSL/TLS", "Terminate certificates on the load balancer\nonly - backends stay simple.", ACCENT_PINK),
]

for i, (title, desc, color) in enumerate(reasons):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(2.4) + Inches(row * 2.6)
    w = Inches(3.7)
    h = Inches(2.2)

    add_shape(slide, x, y, w, h, BG_CARD, 0.05)
    # Color accent bar at top of card
    add_shape(slide, x, y, w, Pt(4), color, 0.0)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.3), w - Inches(0.5), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.85), w - Inches(0.5), Inches(1.2),
                 desc, font_size=13, color=LIGHT_GRAY, line_spacing=1.5)

add_slide_number(slide, 4, TOTAL_SLIDES)

# ===================================================================
# SLIDE 5 — ARCHITECTURE
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "03  ARCHITECTURE")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Architecture & Process Flow", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_PURPLE)

# --- Client box ---
add_shape(slide, Inches(5.2), Inches(2.2), Inches(2.8), Inches(0.7), ACCENT_BLUE, 0.1)
add_text_box(slide, Inches(5.2), Inches(2.25), Inches(2.8), Inches(0.7),
             "Client Request", font_size=16, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# Arrow down
add_shape(slide, Inches(6.45), Inches(2.95), Pt(4), Inches(0.4), ACCENT_BLUE)

# --- VM2 LB box ---
add_shape(slide, Inches(3.8), Inches(3.4), Inches(5.6), Inches(1.5), BG_CARD, 0.05)
add_shape(slide, Inches(3.8), Inches(3.4), Inches(5.6), Pt(4), ACCENT_PURPLE)
add_text_box(slide, Inches(4.0), Inches(3.55), Inches(5.2), Inches(0.4),
             "VM2  -  Nginx Load Balancer  (port 80)", font_size=18, color=ACCENT_PURPLE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.0), Inches(4.0), Inches(5.2), Inches(0.8),
             "upstream: 192.168.1.8:8081  |  192.168.1.8:8082\nAlgorithm: Round-Robin  |  Failover: Automatic",
             font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrows down (two)
add_shape(slide, Inches(5.0), Inches(5.0), Pt(4), Inches(0.35), ACCENT_BLUE)
add_shape(slide, Inches(8.2), Inches(5.0), Pt(4), Inches(0.35), ACCENT_GREEN)

# --- Site A box ---
add_shape(slide, Inches(3.5), Inches(5.4), Inches(3.0), Inches(1.2), BG_CARD, 0.05)
add_shape(slide, Inches(3.5), Inches(5.4), Inches(3.0), Pt(4), ACCENT_BLUE)
add_text_box(slide, Inches(3.5), Inches(5.55), Inches(3.0), Inches(0.4),
             "Site A  (port 8081)", font_size=16, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.5), Inches(5.95), Inches(3.0), Inches(0.5),
             "Primary Web Application\nVM1: 192.168.1.8", font_size=12, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

# --- Site B box ---
add_shape(slide, Inches(6.8), Inches(5.4), Inches(3.0), Inches(1.2), BG_CARD, 0.05)
add_shape(slide, Inches(6.8), Inches(5.4), Inches(3.0), Pt(4), ACCENT_GREEN)
add_text_box(slide, Inches(6.8), Inches(5.55), Inches(3.0), Inches(0.4),
             "Site B  (port 8082)", font_size=16, color=ACCENT_GREEN, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.8), Inches(5.95), Inches(3.0), Inches(0.5),
             "Secondary Web Application\nVM1: 192.168.1.8", font_size=12, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

# --- Dashboard box (left side) ---
add_shape(slide, Inches(0.5), Inches(3.8), Inches(2.8), Inches(1.8), BG_CARD, 0.05)
add_shape(slide, Inches(0.5), Inches(3.8), Inches(2.8), Pt(4), ACCENT_ORANGE)
add_text_box(slide, Inches(0.5), Inches(3.95), Inches(2.8), Inches(0.4),
             "Health Dashboard", font_size=16, color=ACCENT_ORANGE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.5), Inches(4.4), Inches(2.8), Inches(1.0),
             "Port 8080\nAuto-refresh every 10s\nMonitors both backends\nResponse time in ms",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Flow steps on right side
flow_steps = [
    "1.  User accesses http://<VM2_IP>",
    "2.  Nginx selects a backend (round-robin)",
    "3.  Request proxied to VM1 :8081 or :8082",
    "4.  Backend responds through Nginx to client",
    "5.  Dashboard checks health every 10 seconds",
]
add_shape(slide, Inches(10.2), Inches(2.5), Inches(2.8), Inches(4.3), BG_CARD, 0.05)
add_text_box(slide, Inches(10.35), Inches(2.6), Inches(2.5), Inches(0.4),
             "Request Flow", font_size=16, color=ACCENT_PINK, bold=True)
add_multiline_box(slide, Inches(10.35), Inches(3.1), Inches(2.5), Inches(3.5),
                  flow_steps, font_size=11, color=LIGHT_GRAY, line_spacing=1.8)

add_slide_number(slide, 5, TOTAL_SLIDES)

# ===================================================================
# SLIDE 6 — ENVIRONMENT DETAILS
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "04  ENVIRONMENT")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Environment Details", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_ORANGE)

env_items = [
    ("Operating System", "AlmaLinux 10 (RHEL-compatible)", ACCENT_BLUE),
    ("Web Server", "Nginx 1.26", ACCENT_GREEN),
    ("Dynamic Content", "PHP-FPM 8.x", ACCENT_PURPLE),
    ("Firewall", "firewalld (optional)", ACCENT_ORANGE),
    ("SELinux", "Enforcing (custom booleans/ports)", ACCENT_PINK),
    ("VM1 IP Address", "192.168.1.8", ACCENT_BLUE),
    ("VM2 IP Address", "192.168.1.20", ACCENT_GREEN),
    ("Dashboard Port", "8080 (health monitoring)", ACCENT_PURPLE),
]

for i, (label, value, color) in enumerate(env_items):
    col = i % 4
    row = i // 4
    x = Inches(0.8) + Inches(col * 3.1)
    y = Inches(2.4) + Inches(row * 2.3)
    w = Inches(2.8)
    h = Inches(2.0)

    add_shape(slide, x, y, w, h, BG_CARD, 0.05)
    # Top color bar
    add_shape(slide, x, y, w, Pt(4), color, 0.0)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.35), w - Inches(0.4), Inches(0.4),
                 label, font_size=14, color=MID_GRAY)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), w - Inches(0.4), Inches(0.8),
                 value, font_size=18, color=color, bold=True)

add_slide_number(slide, 6, TOTAL_SLIDES)

# ===================================================================
# SLIDE 7 — PHASE 1: VM1 SETUP
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "05  PHASE 1: VM1")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Phase 1: VM1 - Two Websites + Dashboard", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_BLUE)

# Step 1 Card
add_bullet_card(slide, Inches(0.8), Inches(2.3), Inches(5.8), Inches(2.2),
                "Step 1: Install Nginx & PHP-FPM",
                [
                    "  sudo dnf install -y nginx php php-fpm",
                    "  sudo systemctl enable --now nginx php-fpm",
                    "  Provides web serving + dynamic PHP content",
                ],
                title_color=ACCENT_BLUE)

# Step 2 Card
add_bullet_card(slide, Inches(7.0), Inches(2.3), Inches(5.8), Inches(2.2),
                "Step 2: Create Website Directories",
                [
                    "  /var/www/site_a/html   (Site A)",
                    "  /var/www/site_b/html   (Site B)",
                    "  /var/www/dashboard/html (Dashboard)",
                ],
                title_color=ACCENT_GREEN)

# Step 3 Card
add_bullet_card(slide, Inches(0.8), Inches(4.8), Inches(5.8), Inches(2.2),
                "Step 3: Create Sample Pages",
                [
                    "  index.html with welcome message",
                    "  uptime.php for live server uptime",
                    "  Each site has unique styling & identity",
                ],
                title_color=ACCENT_PURPLE)

# Step 4 Card
add_bullet_card(slide, Inches(7.0), Inches(4.8), Inches(5.8), Inches(2.2),
                "Step 4: Configure Nginx Virtual Hosts",
                [
                    "  site_a.conf  -  listen 8081",
                    "  site_b.conf  -  listen 8082",
                    "  dashboard.conf  -  listen 8080 (PHP)",
                ],
                title_color=ACCENT_ORANGE)

add_slide_number(slide, 7, TOTAL_SLIDES)

# ===================================================================
# SLIDE 8 — PHASE 1: CONFIG + SELINUX
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "06  CONFIG & SELINUX")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Phase 1: Configuration & SELinux", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_GREEN)

# Nginx Config code block
add_text_box(slide, Inches(0.8), Inches(2.2), Inches(6), Inches(0.4),
             "Nginx Virtual Host Configuration (site_a.conf)", font_size=15,
             color=ACCENT_BLUE, bold=True)

nginx_code = """server {
    listen 8081;
    root /var/www/site_a/html;
    index index.html index.php;

    location ~ \\.php$ {
        fastcgi_pass unix:/run/php-fpm/www.sock;
        include fastcgi_params;
        fastcgi_param SCRIPT_FILENAME
            $document_root$fastcgi_script_name;
    }
}"""
add_code_block(slide, Inches(0.8), Inches(2.7), Inches(6), Inches(3.2), nginx_code)

# SELinux card
add_text_box(slide, Inches(7.2), Inches(2.2), Inches(5.5), Inches(0.4),
             "SELinux Configuration", font_size=15, color=ACCENT_ORANGE, bold=True)

selinux_code = """# Allow custom ports
sudo semanage port -a -t http_port_t \\
    -p tcp 8081
sudo semanage port -a -t http_port_t \\
    -p tcp 8082
sudo semanage port -a -t http_port_t \\
    -p tcp 8080

# Allow network connections
sudo setsebool -P \\
    httpd_can_network_connect 1"""
add_code_block(slide, Inches(7.2), Inches(2.7), Inches(5.5), Inches(3.2), selinux_code)

# Firewall card
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(5.5), Inches(0.4),
             "Firewall Rules", font_size=15, color=ACCENT_PINK, bold=True)

fw_code = """sudo firewall-cmd --permanent --add-port={8081,8082,8080}/tcp
sudo firewall-cmd --reload"""
add_code_block(slide, Inches(0.8), Inches(6.55), Inches(5.5), Inches(0.7), fw_code)

# Restart card
add_text_box(slide, Inches(7.2), Inches(6.2), Inches(5.5), Inches(0.4),
             "Validate & Restart", font_size=15, color=ACCENT_GREEN, bold=True)

restart_code = """sudo nginx -t && sudo systemctl restart nginx php-fpm"""
add_code_block(slide, Inches(7.2), Inches(6.55), Inches(5.5), Inches(0.7), restart_code)

add_slide_number(slide, 8, TOTAL_SLIDES)

# ===================================================================
# SLIDE 9 — PHASE 2: VM2 LOAD BALANCER
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "07  PHASE 2: VM2")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Phase 2: VM2 - Nginx Load Balancer", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_PURPLE)

# Left: LB config
add_text_box(slide, Inches(0.8), Inches(2.2), Inches(6), Inches(0.4),
             "Load Balancer Config  (/etc/nginx/conf.d/load_balancer.conf)",
             font_size=15, color=ACCENT_PURPLE, bold=True)

lb_code = """upstream backend_servers {
    server 192.168.1.8:8081;
    server 192.168.1.8:8082;
}

server {
    listen 80;

    location / {
        proxy_pass http://backend_servers;
        proxy_set_header Host $host;
        proxy_set_header X-Real-IP $remote_addr;
    }

    location /dashboard {
        proxy_pass http://192.168.1.8:8080;
    }
}"""
add_code_block(slide, Inches(0.8), Inches(2.7), Inches(6.2), Inches(4.2), lb_code)

# Right side: steps
steps_data = [
    ("Install Nginx", "sudo dnf install -y nginx", ACCENT_BLUE),
    ("SELinux Boolean", "sudo setsebool -P httpd_can_network_connect 1", ACCENT_GREEN),
    ("Open Firewall", "sudo firewall-cmd --permanent --add-service=http\nsudo firewall-cmd --reload", ACCENT_ORANGE),
    ("Start Nginx", "sudo systemctl enable --now nginx", ACCENT_PINK),
]

for i, (title, cmd, color) in enumerate(steps_data):
    y = Inches(2.3) + Inches(i * 1.25)
    add_shape(slide, Inches(7.4), y, Inches(5.4), Inches(1.1), BG_CARD, 0.05)
    add_shape(slide, Inches(7.4), y, Pt(5), Inches(1.1), color, 0.0)
    add_text_box(slide, Inches(7.7), y + Inches(0.1), Inches(5), Inches(0.35),
                 f"Step {i+1}: {title}", font_size=15, color=color, bold=True)
    add_text_box(slide, Inches(7.7), y + Inches(0.5), Inches(5), Inches(0.5),
                 cmd, font_size=11, color=ACCENT_GREEN, font_name="Courier New")

add_slide_number(slide, 9, TOTAL_SLIDES)

# ===================================================================
# SLIDE 10 — DEMO: SITE A & SITE B (screenshots)
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "08  DEMO: WEBSITES")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Demo: Site A & Site B", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_ORANGE)

# Site A image
img_a = os.path.join(IMG_DIR, "resized", "Websites-b.png")  # This is actually Site A
img_b = os.path.join(IMG_DIR, "resized", "Websites-a.png")  # This is actually Site B

if os.path.exists(img_a):
    add_shape(slide, Inches(0.6), Inches(2.3), Inches(6.1), Inches(3.65), BG_CARD, 0.03)
    slide.shapes.add_picture(img_a, Inches(0.7), Inches(2.4), Inches(5.9), Inches(3.45))
    add_text_box(slide, Inches(0.7), Inches(6.1), Inches(5.9), Inches(0.5),
                 "Site A  -  Primary Web Application  (port 8081)",
                 font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.7), Inches(6.5), Inches(5.9), Inches(0.5),
                 "Server IP: 192.168.1.8  |  Status: Healthy & Running  |  Uptime: 49 min",
                 font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

if os.path.exists(img_b):
    add_shape(slide, Inches(6.9), Inches(2.3), Inches(6.1), Inches(3.65), BG_CARD, 0.03)
    slide.shapes.add_picture(img_b, Inches(7.0), Inches(2.4), Inches(5.9), Inches(3.45))
    add_text_box(slide, Inches(7.0), Inches(6.1), Inches(5.9), Inches(0.5),
                 "Site B  -  Secondary Web Application  (port 8082)",
                 font_size=14, color=ACCENT_PINK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.0), Inches(6.5), Inches(5.9), Inches(0.5),
                 "Server IP: 192.168.1.8  |  Status: Healthy & Running  |  Uptime: 49 min",
                 font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 10, TOTAL_SLIDES)

# ===================================================================
# SLIDE 11 — DEMO: HEALTH DASHBOARD
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "09  HEALTH DASHBOARD")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Demo: Health Dashboard", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_PINK)

img_dash = os.path.join(IMG_DIR, "resized", "HealthDashboard.png")
if os.path.exists(img_dash):
    add_shape(slide, Inches(0.6), Inches(2.3), Inches(8.4), Inches(4.7), BG_CARD, 0.03)
    slide.shapes.add_picture(img_dash, Inches(0.7), Inches(2.4), Inches(8.2), Inches(4.5))

# Info cards on right
info_cards = [
    ("Auto-Refresh", "Updates every 10\nseconds automatically", ACCENT_BLUE),
    ("Response Time", "Measures latency in\nmilliseconds per backend", ACCENT_GREEN),
    ("Status Monitor", "Green = UP, Red = DOWN\nReal-time health check", ACCENT_ORANGE),
    ("Dual Backend", "Monitors both Site A\n(:8081) and Site B (:8082)", ACCENT_PINK),
]

for i, (title, desc, color) in enumerate(info_cards):
    y = Inches(2.3) + Inches(i * 1.2)
    add_shape(slide, Inches(9.3), y, Inches(3.6), Inches(1.05), BG_CARD, 0.05)
    add_shape(slide, Inches(9.3), y, Pt(4), Inches(1.05), color, 0.0)
    add_text_box(slide, Inches(9.6), y + Inches(0.08), Inches(3.2), Inches(0.3),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(9.6), y + Inches(0.45), Inches(3.2), Inches(0.55),
                 desc, font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 11, TOTAL_SLIDES)

# ===================================================================
# SLIDE 12 — DEMO: NGINX LB IN ACTION
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "10  LB IN ACTION")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Demo: Nginx Load Balancer in Action", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_BLUE)

img_nginx = os.path.join(IMG_DIR, "resized", "Nginxstatus.png")
if os.path.exists(img_nginx):
    add_shape(slide, Inches(0.6), Inches(2.3), Inches(8.4), Inches(4.7), BG_CARD, 0.03)
    slide.shapes.add_picture(img_nginx, Inches(0.7), Inches(2.4), Inches(8.2), Inches(4.5))

# Highlight cards on right
highlights = [
    ("Round-Robin Proof", "curl shows alternating\nresponses: Site A then\nSite B then Site A...", ACCENT_BLUE),
    ("Access Logs", "tail -f shows requests\nbeing distributed across\nboth backends", ACCENT_GREEN),
    ("Service Status", "systemctl status nginx\nshows active (running)\non VM2", ACCENT_PURPLE),
    ("HTTP 200 OK", "All responses return\nHTTP 200 - backends\nare healthy", ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(highlights):
    y = Inches(2.3) + Inches(i * 1.2)
    add_shape(slide, Inches(9.3), y, Inches(3.6), Inches(1.05), BG_CARD, 0.05)
    add_shape(slide, Inches(9.3), y, Pt(4), Inches(1.05), color, 0.0)
    add_text_box(slide, Inches(9.6), y + Inches(0.08), Inches(3.2), Inches(0.3),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(9.6), y + Inches(0.45), Inches(3.2), Inches(0.55),
                 desc, font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 12, TOTAL_SLIDES)

# ===================================================================
# SLIDE 13 — VERIFICATION & TESTING
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "11  TESTING")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Verification & Testing", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_GREEN)

tests = [
    ("Site A Direct", "curl http://192.168.1.8:8081", "Welcome message or styled page", ACCENT_BLUE),
    ("Site B Direct", "curl http://192.168.1.8:8082", "Welcome message or styled page", ACCENT_GREEN),
    ("Dashboard", "curl http://192.168.1.8:8080", "Health status JSON or HTML", ACCENT_PURPLE),
    ("Load Balancer", "curl http://192.168.1.20 (x3)", "Alternating: Site A, B, A...", ACCENT_ORANGE),
    ("Dashboard via LB", "curl http://192.168.1.20/dashboard", "Dashboard HTML page", ACCENT_PINK),
    ("Failure Test", "sudo systemctl stop nginx (VM1)", "VM2 stops routing to failed backend", RGBColor(0xEF, 0x44, 0x44)),
]

for i, (test_name, command, expected, color) in enumerate(tests):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(2.3) + Inches(row * 1.6)
    w = Inches(5.8)
    h = Inches(1.4)

    add_shape(slide, x, y, w, h, BG_CARD, 0.05)
    add_shape(slide, x, y, w, Pt(3), color, 0.0)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.3),
                 test_name, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), w - Inches(0.4), Inches(0.3),
                 command, font_size=12, color=ACCENT_GREEN, font_name="Courier New")
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), w - Inches(0.4), Inches(0.35),
                 f"Expected: {expected}", font_size=12, color=LIGHT_GRAY)

add_slide_number(slide, 13, TOTAL_SLIDES)

# ===================================================================
# SLIDE 14 — CHALLENGES & SOLUTIONS
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "12  CHALLENGES")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Challenges & Solutions", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_PURPLE)

challenges = [
    ("SELinux Blocking Ports", "Used semanage port -a to add\ncustom HTTP ports", ACCENT_BLUE),
    ("PHP-FPM Network Error", "Enabled httpd_can_network_connect\nSELinux boolean", ACCENT_GREEN),
    ("Port 80 Conflict", "Stopped & removed HAProxy that\nwas occupying the port", ACCENT_PURPLE),
    ("Nginx Syntax Errors", "Replaced config & validated with\nnginx -t before restart", ACCENT_ORANGE),
    ("Emoji Display Issues", "Added <meta charset=\"UTF-8\">\nto all HTML pages", ACCENT_PINK),
    ("PHP Socket Permission", "Fixed via SELinux boolean +\nPHP-FPM socket permissions", ACCENT_BLUE),
]

for i, (challenge, solution, color) in enumerate(challenges):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(2.3) + Inches(row * 2.5)
    w = Inches(3.7)
    h = Inches(2.2)

    add_shape(slide, x, y, w, h, BG_CARD, 0.05)
    add_shape(slide, x, y, w, Pt(4), color, 0.0)
    # Challenge
    add_text_box(slide, x + Inches(0.2), y + Inches(0.25), w - Inches(0.4), Inches(0.35),
                 "CHALLENGE", font_size=10, color=MID_GRAY, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), w - Inches(0.4), Inches(0.4),
                 challenge, font_size=15, color=color, bold=True)
    # Solution
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), w - Inches(0.4), Inches(0.3),
                 "SOLUTION", font_size=10, color=MID_GRAY, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.4), w - Inches(0.4), Inches(0.7),
                 solution, font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 14, TOTAL_SLIDES)

# ===================================================================
# SLIDE 15 — EXTENDING THE INFRASTRUCTURE
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_section_badge(slide, "13  EXTENDING")

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "Extending the Infrastructure", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.8), Inches(2.5), ACCENT_ORANGE)

extensions = [
    ("Add More Backends",
     "Simply add new server lines in the\nupstream block. No client-side changes\nneeded - instant horizontal scaling.",
     ACCENT_BLUE),
    ("Load Balancing Algorithms",
     "Switch from round-robin to:\n  least_conn;  (fewest connections)\n  ip_hash;  (sticky sessions)\n  random;  (randomized)",
     ACCENT_GREEN),
    ("Enable SSL/TLS",
     "Obtain a certificate (Let's Encrypt) and\nadd listen 443 ssl; on VM2.\nBackends communicate over plain HTTP.",
     ACCENT_PURPLE),
    ("Session Persistence",
     "Use sticky cookie (Nginx Plus) or\nip_hash for session affinity.\nKeeps user on same backend.",
     ACCENT_ORANGE),
    ("Logging & Monitoring",
     "Access logs show which upstream handled\neach request. Integrate with Grafana or\nPrometheus for visualization.",
     ACCENT_PINK),
]

for i, (title, desc, color) in enumerate(extensions):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(2.4) + Inches(row * 2.6)
    w = Inches(3.7)
    h = Inches(2.3)

    add_shape(slide, x, y, w, h, BG_CARD, 0.05)
    add_shape(slide, x, y, w, Pt(4), color, 0.0)

    add_text_box(slide, x + Inches(0.25), y + Inches(0.3), w - Inches(0.5), Inches(0.4),
                 title, font_size=17, color=color, bold=True)
    add_accent_line(slide, x + Inches(0.25), y + Inches(0.7), Inches(1.2), color)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.85), w - Inches(0.5), Inches(1.3),
                 desc, font_size=12, color=LIGHT_GRAY, line_spacing=1.5)

add_slide_number(slide, 15, TOTAL_SLIDES)

# ===================================================================
# SLIDE 16 — CONCLUSION
# ===================================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# Decorative shapes
add_shape(slide, Inches(-1), Inches(-1), Inches(6), Inches(6),
          RGBColor(0x38, 0xBD, 0xF8), 0.5).fill.fore_color.rgb = RGBColor(0x14, 0x28, 0x50)
add_shape(slide, Inches(9), Inches(3), Inches(6), Inches(6),
          RGBColor(0x38, 0xBD, 0xF8), 0.5).fill.fore_color.rgb = RGBColor(0x14, 0x28, 0x50)

add_section_badge(slide, "14  CONCLUSION")

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(0.7),
             "Conclusion", font_size=42, color=WHITE, bold=True)
add_accent_line(slide, Inches(1.5), Inches(2.5), Inches(3), ACCENT_PINK)

conclusion_points = [
    "Production-ready, scalable, and monitored web infrastructure",
    "Built entirely with open-source tools (Nginx + PHP-FPM)",
    "Load balancer ensures high availability with automatic failover",
    "Health dashboard provides real-time visibility into backend status",
    "SELinux enforcing mode maintained for enterprise-grade security",
    "All common challenges addressed with practical solutions",
]

for i, point in enumerate(conclusion_points):
    y = Inches(3.0) + Inches(i * 0.55)
    colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_PURPLE, ACCENT_ORANGE, ACCENT_PINK, ACCENT_BLUE]
    # Bullet dot
    add_shape(slide, Inches(1.8), y + Inches(0.08), Inches(0.15), Inches(0.15), colors[i], 0.5)
    add_text_box(slide, Inches(2.2), y, Inches(9), Inches(0.45),
                 point, font_size=18, color=LIGHT_GRAY)

# Footer info
add_accent_line(slide, Inches(1.5), Inches(6.2), Inches(10), MID_GRAY)
add_text_box(slide, Inches(1.5), Inches(6.4), Inches(10), Inches(0.5),
             "AlmaLinux 10  |  Nginx 1.26  |  PHP-FPM 8.x  |  April 2026",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3.5), Inches(6.8), Inches(6), Inches(0.5),
             "Thank You!", font_size=32, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 16, TOTAL_SLIDES)

# ===================================================================
# SAVE
# ===================================================================
output_path = os.path.join(IMG_DIR, "LoadBalancer_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
